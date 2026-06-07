# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(BASE, "slides")

# ---- 밝은 컬러 팔레트 ----
BG = RGBColor(0xF7, 0xF8, 0xFC)        # 거의 흰색, 살짝 푸른 톤
PANEL = RGBColor(0xFF, 0xFF, 0xFF)     # 카드 배경(흰색)
PANEL_BORDER = RGBColor(0xDD, 0xE2, 0xF0)
ACCENT = RGBColor(0xE2, 0x8A, 0x12)    # 주황/금색 (태양 느낌)
ACCENT2 = RGBColor(0x2F, 0x6F, 0xE0)   # 파란색
TEXT = RGBColor(0x23, 0x27, 0x3A)
DIM = RGBColor(0x6B, 0x72, 0x8E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, text, size, color=TEXT,
             bold=False, align=PP_ALIGN.LEFT, font="Malgun Gothic", line_spacing=1.0,
             anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_eyebrow_title(slide, eyebrow, title):
    add_text(slide, Inches(0.7), Inches(0.45), Inches(8), Inches(0.4),
             eyebrow, 14, ACCENT2, bold=True, font="Consolas")
    add_text(slide, Inches(0.7), Inches(0.85), Inches(11.8), Inches(0.9),
             title, 30, ACCENT, bold=True)


def add_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = PANEL
    card.line.color.rgb = PANEL_BORDER
    card.line.width = Pt(1)
    card.shadow.inherit = False
    return card


def add_card_with_text(slide, left, top, width, height, heading, body, heading_size=16, body_size=13):
    add_card(slide, left, top, width, height)
    pad = Inches(0.28)
    add_text(slide, left + pad, top + Inches(0.2), width - pad * 2, Inches(0.4),
             heading, heading_size, ACCENT2, bold=True)
    add_text(slide, left + pad, top + Inches(0.65), width - pad * 2, height - Inches(0.85),
             body, body_size, DIM, line_spacing=1.18)


def add_pagenum(slide, n, total):
    add_text(slide, SW - Inches(1.6), SH - Inches(0.55), Inches(1.2), Inches(0.4),
             "{} / {}".format(n, total), 12, DIM, align=PP_ALIGN.RIGHT, font="Consolas")


TOTAL = 9
PAGE = [0]


def next_page():
    PAGE[0] += 1
    return PAGE[0]


# =====================================================================
# 0. 표지
# =====================================================================
s = add_slide()
add_text(s, Inches(0), Inches(2.5), SW, Inches(0.5),
         "COMPUTER GRAPHICS TERM PROJECT", 16, ACCENT2, bold=True,
         align=PP_ALIGN.CENTER, font="Consolas")
add_text(s, Inches(0), Inches(3.05), SW, Inches(1.3),
         "🪐 태양계 관측 시뮬레이터", 44, ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(2), Inches(4.35), Inches(9.333), Inches(1.2),
         "WebGL2 기반 실시간 3D 태양계 비행 관측 시뮬레이터\n"
         "절차적 지오메트리 · 커스텀 셰이더 · 프레임버퍼 피킹으로 만든 1인 프로젝트",
         18, DIM, align=PP_ALIGN.CENTER, line_spacing=1.3)
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ① 주제
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "① TOPIC", "주제 — 태양계 관측 시뮬레이터")
add_text(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(2.0),
         "우주선 시점의 자유 비행 카메라로 태양계 안을 직접 날아다니며, 태양과 8개 행성을 "
         "가까이서 관찰하고 정보를 확인하는 실시간 3D WebGL2 애플리케이션입니다.\n\n"
         "시간 배속 조절 · 초점 추적 · 실제/압축 비율 전환 등 '관측'이라는 행위 자체를 "
         "능동적인 상호작용으로 설계했습니다.",
         17, TEXT, line_spacing=1.35)

tags = ["WebGL2", "실시간 3D 렌더링", "절차적 생성", "자유비행 카메라", "1인 개발"]
tx = Inches(0.7)
for tag in tags:
    w = Inches(0.35 + 0.16 * len(tag))
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tx, Inches(4.5), w, Inches(0.55))
    chip.adjustments[0] = 0.5
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(0xE9, 0xF1, 0xFD)
    chip.line.color.rgb = ACCENT2
    chip.line.width = Pt(1)
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tag
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ACCENT2
    tx += w + Inches(0.25)
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ② 내용 (제작 동기, 스토리)
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "② BACKGROUND & STORY", "내용 — 제작 동기와 스토리")
cw = Inches(5.85)
ch = Inches(2.35)
add_card_with_text(s, Inches(0.7), Inches(2.0), cw, ch,
                   "처음 계획: \"인터랙티브 3D 갤러리\"",
                   "OBJ 모델 로딩과 피킹 기능을 살려 전시관 갤러리를 기획·완성했지만, "
                   "완성 후 보니 흔한 소재라 \"식상하다\"는 느낌이 강하게 들었습니다.")
add_card_with_text(s, Inches(6.78), Inches(2.0), cw, ch,
                   "방향 전환 — 태양계로",
                   "갤러리는 보존해 두고, 같은 핵심 기술(변환 행렬·조명·텍스처·피킹)을 "
                   "누구에게나 친숙하면서 시각적으로 화려한 태양계 소재에 새롭게 적용했습니다.")
add_card_with_text(s, Inches(0.7), Inches(4.55), Inches(11.93), Inches(2.05),
                   "강의 자료에 없던 것을 직접 채워나간 과정",
                   "구 절차적 생성, 행성 텍스처, 궤도 운동 같은 요소는 강의 자료에 없었지만, "
                   "오히려 그것이 \"배운 개념을 새 문제에 적용해보는\" 도전이 되었고, "
                   "결국 끝까지 완성해낸 결과물입니다.")
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ③ 주요 구현 기술
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "③ IMPLEMENTATION", "주요 구현 기술")
cw = Inches(5.85)
ch = Inches(2.35)
items = [
    ("1) 절차적 지오메트리",
     "외부 3D 모델 없이 generateSphere(위도·경도 분할), generateRing, "
     "generateCircleOutline, generateStarField로 행성·고리·궤도선·별을 모두 직접 생성했습니다."),
    ("2) 셰이더 3종 — 2개를 직접 변형/신규 작성",
     "행성 셰이더: 평행광 모델 → 점광원 모델로 변형 (lightDir = normalize(sunPos − worldPos))\n"
     "태양 셰이더(신규): emissive + 노이즈 + 펄싱 코로나\n"
     "uniform_color: 피킹·별·궤도선·하이라이트에 다용도 재사용"),
    ("3) 프레임버퍼 색상 피킹",
     "클릭 시 별도 프레임버퍼에 각 천체를 고유 ID 색으로 렌더링한 뒤 픽셀을 읽어 "
     "선택된 천체를 정확히 판별합니다."),
    ("4) 1인칭 자유비행 카메라",
     "yaw/pitch 기반 드래그 회전, 외적으로 구한 right 벡터로 WASD 스트레이프 이동, "
     "행성은 translate·rotateY·scale 계층적 변환으로 배치했습니다."),
]
positions = [(Inches(0.7), Inches(2.0)), (Inches(6.78), Inches(2.0)),
             (Inches(0.7), Inches(4.55)), (Inches(6.78), Inches(4.55))]
for (heading, body), (l, t) in zip(items, positions):
    add_card_with_text(s, l, t, cw, ch, heading, body, body_size=12.5)
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ④ 차별성
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "④ WHAT MAKES IT STAND OUT", "차별성 — 특징 및 자랑할 만한 점")
cw = Inches(3.79)
ch = Inches(2.35)
items = [
    ("🔭 초점 추적 카메라", "행성 클릭 시 카메라가 부드럽게 자동으로 따라가며 비춥니다. 드래그하면 즉시 추적이 해제됩니다."),
    ("📏 실제/압축 비율 전환", "R 키로 '관측하기 좋은 비율'과 '실제 천문학적 비율'을 즉시 비교할 수 있습니다."),
    ("⏯ 시간 배속 / 일시정지", "+/− 로 0.25~8배속 조절, Space로 일시정지 — 배속을 기억했다 복원합니다."),
    ("🎯 정보 패널 + 하이라이트 링", "천체를 클릭하면 설명 패널과 펄싱 하이라이트 링이 나타나 무엇을 보고 있는지 직관적으로 알 수 있습니다."),
    ("🖼 메인화면 / 조작 방법", "타이틀 화면에서 '시작'·'조작 방법'을 선택하며, 상태 머신으로 화면 전환을 깔끔히 분리했습니다."),
    ("🌌 외부 모델 없이 절차적 완성", "행성·고리·궤도선·별 배경까지 모든 지오메트리를 코드로 직접 생성한 점이 핵심 자부심입니다."),
]
xs = [Inches(0.7), Inches(4.78), Inches(8.86)]
ys = [Inches(2.0), Inches(4.55)]
for idx, (heading, body) in enumerate(items):
    l = xs[idx % 3]
    t = ys[idx // 3]
    add_card_with_text(s, l, t, cw, ch, heading, body, heading_size=14, body_size=12)
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ⑤ 사용자 인터페이스
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "⑤ USER INTERFACE", "사용자 인터페이스 — 조작법")
add_card_with_text(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.6),
                   "화면 흐름",
                   "실행 → 타이틀 화면('시작'/'조작 방법') → '시작'으로 바로 진입, "
                   "'조작 방법'으로 단축키 안내 화면 이동.\n\n"
                   "시뮬레이션 중에도 화면 좌상단에 핵심 단축키가 항상 표시되어 "
                   "언제든 조작법을 다시 확인할 수 있습니다.\n\n"
                   "타이틀/도움말 화면에서는 키보드·클릭 입력이 시뮬레이션에 영향을 "
                   "주지 않도록 상태 머신으로 깔끔하게 분리했습니다.",
                   body_size=13)

card2 = add_card(s, Inches(6.78), Inches(2.0), Inches(5.85), Inches(4.6))
add_text(s, Inches(7.06), Inches(2.2), Inches(5.3), Inches(0.4), "조작법 한눈에 보기", 16, ACCENT2, bold=True)
controls = [
    ("W A S D", "이동 (전후좌우)"),
    ("드래그", "시점 회전"),
    ("휠", "확대 / 축소"),
    ("클릭", "천체 선택 → 정보 패널 + 자동 추적"),
    ("+ / −", "공전 속도 배율 조절"),
    ("Space", "일시정지 / 재생"),
    ("R", "실제 ↔ 압축 비율 전환"),
]
cy = Inches(2.78)
row_h = Inches(0.55)
for key, desc in controls:
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.06), cy, Inches(1.55), Inches(0.42))
    chip.adjustments[0] = 0.25
    chip.fill.solid()
    chip.fill.fore_color.rgb = RGBColor(0xFB, 0xF1, 0xDD)
    chip.line.color.rgb = ACCENT
    chip.line.width = Pt(0.75)
    chip.shadow.inherit = False
    tf = chip.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = key
    run.font.size = Pt(12.5)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = "Consolas"
    add_text(s, Inches(8.78), cy + Inches(0.06), Inches(3.7), Inches(0.4), desc, 12.5, DIM)
    cy += row_h
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ⑥ 향후 개선점
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "⑥ FUTURE IMPROVEMENTS", "향후 개선점")
add_card_with_text(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.6),
                   "✋ 어려웠던 점 — \"비율의 딜레마\"",
                   "실제 비율로는 행성이 점처럼 보이고, 압축하면 정확하지 않습니다. "
                   "둘을 동시에 만족하는 단일 스케일은 없다고 판단해, 두 모드를 모두 구현해 "
                   "사용자가 직접 비교하는 절충안을 택했습니다.",
                   body_size=13.5)
card2 = add_card(s, Inches(6.78), Inches(2.0), Inches(5.85), Inches(4.6))
add_text(s, Inches(7.06), Inches(2.2), Inches(5.3), Inches(0.4), "🚀 시간이 더 있다면", 16, ACCENT2, bold=True)
future = [
    "노멀 매핑 — 행성 표면에 입체적인 굴곡 표현",
    "소행성대 파티클 — 화성·목성 사이를 더 풍부하게",
    "그림자 / 일식 — 점광원 기반 섀도우 매핑",
    "위성(달) 추가 — 행성 주위를 도는 계층적 공전",
]
fy = Inches(2.85)
for item in future:
    add_text(s, Inches(7.06), fy, Inches(5.3), Inches(0.8), "•  " + item, 13.5, DIM, line_spacing=1.2)
    fy += Inches(0.85)
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ⑦ 스크린샷
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "⑦ SCREENSHOTS", "스크린샷")
shots = [
    ("slide_title.png", "① 메인화면", "제목·설명과 '시작'/'조작 방법' 메뉴"),
    ("slide_selected.png", "② 천체 선택 + 추적", "토성 클릭 → 정보 패널·하이라이트 링·자동 추적"),
    ("slide_realscale.png", "③ 실제 비율 모드", "R 키로 전환한 압도적인 실제 스케일"),
]
iw = Inches(3.85)
ih = Inches(2.55)
gap = Inches(0.34)
total_w = iw * 3 + gap * 2
left0 = (SW - total_w) // 2
for i, (fname, heading, desc) in enumerate(shots):
    l = left0 + i * (iw + gap)
    t = Inches(2.05)
    frame = add_card(s, l, t, iw, ih + Inches(1.15))
    img_path = os.path.join(SLIDES_DIR, fname)
    pic = s.shapes.add_picture(img_path, l + Inches(0.12), t + Inches(0.12), width=iw - Inches(0.24), height=ih)
    add_text(s, l + Inches(0.22), t + ih + Inches(0.22), iw - Inches(0.4), Inches(0.4), heading, 14, ACCENT2, bold=True)
    add_text(s, l + Inches(0.22), t + ih + Inches(0.62), iw - Inches(0.4), Inches(0.7), desc, 11.5, DIM, line_spacing=1.15)
add_pagenum(s, next_page(), TOTAL)

# =====================================================================
# ⑧ 에셋 출처
# =====================================================================
s = add_slide()
add_eyebrow_title(s, "⑧ ASSET CREDITS", "에셋 출처")
add_card_with_text(s, Inches(0.7), Inches(2.0), Inches(11.93), Inches(2.05),
                   "행성 / 고리 텍스처",
                   "Solar System Scope (solarsystemscope.com/textures)의 NASA 이미지 기반 "
                   "텍스처를, CC BY 4.0 (creativecommons.org/licenses/by/4.0) 라이선스에 따라 "
                   "사용했습니다.",
                   body_size=14)
add_card_with_text(s, Inches(0.7), Inches(4.3), Inches(11.93), Inches(2.3),
                   "그 외 — 전부 직접 제작",
                   "지오메트리 · 셰이더(GLSL) · UI(HTML/CSS) · 로직(JavaScript/WebGL2) 모두 "
                   "직접 작성했으며, 외부 3D 모델·사운드·폰트 등의 추가 에셋은 사용하지 않았습니다.\n"
                   "라이선스 전문은 프로젝트 폴더의 textures/LICENSE_NOTICE.txt 파일에 있습니다.",
                   body_size=14)
add_pagenum(s, next_page(), TOTAL)

out_path = os.path.join(BASE, "presentation.pptx")
prs.save(out_path)
print("Saved:", out_path)
